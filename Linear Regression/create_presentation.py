from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_ppt():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Student Score Prediction"
    subtitle.text = "Implementation of Linear Regression\nUsing Python & Scikit-Learn"

    # Slide 2: What is Linear Regression?
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "What is Linear Regression?"
    content = slide.placeholders[1]
    content.text = ("Linear Regression is a statistical method that models the relationship "
                    "between a dependent variable (Score) and one or more independent variables (Hours).\n\n"
                    "Goal: To find the 'Line of Best Fit' that minimizes the error between actual and predicted values.")

    # Slide 3: The Dataset
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Dataset (student_scores.csv)"
    content = slide.placeholders[1]
    content.text = "A snippet of the data used for training:"
    
    # Adding a table for data snippet
    rows, cols = 6, 2
    left = Inches(2.0)
    top = Inches(2.5)
    width = Inches(4.0)
    height = Inches(0.8)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.0)
    
    table.cell(0, 0).text = 'Hours'
    table.cell(0, 1).text = 'Scores'
    
    sample_data = [('2.5', '21'), ('5.1', '47'), ('3.2', '27'), ('8.5', '75'), ('3.5', '30')]
    for i, (h, s) in enumerate(sample_data):
        table.cell(i+1, 0).text = h
        table.cell(i+1, 1).text = s

    # Slide 4: Model Training Code
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Model Training Code Snippet"
    content = slide.placeholders[1]
    
    code_text = (
        "from sklearn.linear_model import LinearRegression\n"
        "regressor = LinearRegression()\n"
        "regressor.fit(X_train, y_train)\n\n"
        "# Predict for a new value\n"
        "prediction = regressor.predict([[9.25]])"
    )
    
    # Adding a text box for code with mono-font
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(3))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = code_text
    p.font.name = 'Courier New'
    p.font.size = Pt(18)

    # Slide 5: Data Visualization (Graphs)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Visualizing the Correlation"
    
    # Add initial plot
    if os.path.exists('initial_plot.png'):
        slide.shapes.add_picture('initial_plot.png', Inches(0.5), Inches(2), height=Inches(4))
    
    # Slide 6: Regression Result (Graphs)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "The Regression Line"
    
    if os.path.exists('regression_result.png'):
        slide.shapes.add_picture('regression_result.png', Inches(0.5), Inches(2), height=Inches(4))

    # Slide 7: Results and Accuracy
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Project Results (With Outliers)"
    content = slide.placeholders[1]
    content.text = (
        "Accuracy (R-Squared Score): -0.057 (Negative Accuracy!)\n"
        "Mean Absolute Error: 14.71\n\n"
        "Impact of Outliers:\n"
        "Input: 9.25 Hours of Study/Day\n"
        "Result: 71.58% Predicted Score\n\n"
        "Lesson: Outliers can significantly skew a Linear Regression model."
    )

    # Save the presentation
    prs.save('Student_Score_Project.pptx')
    print("PPT generated successfully as 'Student_Score_Project.pptx'")

import os
if __name__ == "__main__":
    create_ppt()
